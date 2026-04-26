"""Shared layout helpers for PPT builders.

Used by `build_review_ppt.py` and `build_paper_ppt_ko.py`.
16:9 layout (13.333" × 7.5"), paper-blue accent, paper-grade typography.
"""

from __future__ import annotations

from pathlib import Path

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ---- Layout constants ----

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ACCENT = RGBColor(0x2E, 0x4A, 0x7F)   # paper-ish blue
GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
GRAY_MID = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_text_box(slide, x, y, w, h, text, *, size=18, bold=False,
                 color=GRAY_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name=None):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if font_name:
        r.font.name = font_name
    return tx


def add_bullets(slide, x, y, w, h, items, *,
                size=16, bullet_color=GRAY_DARK, leading=1.15, font_name=None):
    """Add a list of bullet items. Item can be a string or (text, sub_size)."""
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, sub_size = item
        else:
            text, sub_size = item, size
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = leading
        r = p.add_run()
        r.text = "• " + text
        r.font.size = Pt(sub_size)
        r.font.color.rgb = bullet_color
        if font_name:
            r.font.name = font_name
    return tx


def add_figure(slide, path: Path, x, y, w=None, h=None):
    if not path.exists():
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                      w or Inches(4), h or Inches(3))
        rect.fill.solid()
        rect.fill.fore_color.rgb = GRAY_LIGHT
        rect.line.color.rgb = GRAY_MID
        tf = rect.text_frame
        tf.text = f"[missing: {path.name}]"
        return rect
    if w is None and h is None:
        return slide.shapes.add_picture(str(path), x, y)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def add_caption(slide, x, y, w, text, *, font_name=None):
    """Italic caption under a figure."""
    tx = slide.shapes.add_textbox(x, y, w, Inches(0.4))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.italic = True
    r.font.color.rgb = GRAY_MID
    if font_name:
        r.font.name = font_name


def add_title_bar(slide, title, *, subtitle=None, font_name=None):
    """Top title bar. Returns the y position below the bar for content."""
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(1.0))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    add_text_box(slide, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.6),
                 title, size=26, bold=True, color=ACCENT,
                 anchor=MSO_ANCHOR.MIDDLE, font_name=font_name)
    if subtitle:
        add_text_box(slide, Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.4),
                     subtitle, size=14, color=GRAY_MID,
                     anchor=MSO_ANCHOR.MIDDLE, font_name=font_name)
        return Inches(1.2)
    return Inches(0.95)


def add_footer(slide, idx, total, footer_text, *, font_name=None):
    add_text_box(slide, Inches(0.3), Inches(7.1), Inches(8.0), Inches(0.3),
                 footer_text, size=9, color=GRAY_LIGHT, font_name=font_name)
    add_text_box(slide, Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3),
                 f"{idx} / {total}", size=9, color=GRAY_LIGHT,
                 align=PP_ALIGN.RIGHT, font_name=font_name)


def add_table(slide, x, y, w, h, rows, *,
              header_color=ACCENT, font_size=12, font_name=None):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = table_shape.table
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    if font_name:
                        run.font.name = font_name
                    if i == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    RGBColor(0xF6, 0xF8, 0xFC) if i % 2 == 0
                    else RGBColor(0xFF, 0xFF, 0xFF))
    return table_shape


def new_slide(prs):
    """Return a new blank slide on `prs`."""
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)
