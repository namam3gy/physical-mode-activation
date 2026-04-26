"""Build a review PPT for advisor / collaborators summarizing the project.

Generates `docs/review_ppt/physical_mode_review_<ts>.pptx` with ~36 slides
covering motivation, methods, M1-M9 milestones, §4.X extensions, and
open backlog. Figures from `docs/figures/` are embedded throughout.

Usage:
    uv run python scripts/build_review_ppt.py
"""

from __future__ import annotations

import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


PROJECT_ROOT = Path(__file__).resolve().parents[1]
FIG_DIR = PROJECT_ROOT / "docs" / "figures"


# ---- Layout helpers ----

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ACCENT = RGBColor(0x2E, 0x4A, 0x7F)   # paper-ish blue
GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
GRAY_MID = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_text_box(slide, x, y, w, h, text, *, size=18, bold=False,
                 color=GRAY_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tx


def add_bullets(slide, x, y, w, h, items: list[str | tuple], *,
                size=16, bullet_color=GRAY_DARK, leading=1.15):
    """Add a list of bullet items. Item can be a string or (text, sub_size)."""
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, sub_size = item
        else:
            text, sub_size = item, size
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = leading
        r = p.add_run()
        r.text = "• " + text
        r.font.size = Pt(sub_size)
        r.font.color.rgb = bullet_color
    return tx


def add_figure(slide, path: Path, x, y, w=None, h=None):
    if not path.exists():
        # Placeholder rectangle
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w or Inches(4), h or Inches(3))
        rect.fill.solid()
        rect.fill.fore_color.rgb = GRAY_LIGHT
        rect.line.color.rgb = GRAY_MID
        tf = rect.text_frame
        tf.text = f"[missing: {path.name}]"
        return rect
    # Actual figure
    if w is None and h is None:
        return slide.shapes.add_picture(str(path), x, y)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def add_caption(slide, x, y, w, text):
    """Italic caption under a figure."""
    tx = slide.shapes.add_textbox(x, y, w, Inches(0.4))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.italic = True
    r.font.color.rgb = GRAY_MID


def add_title_bar(slide, title: str, *, subtitle: str | None = None):
    """Top title bar. Returns the y position below the bar for content."""
    # Accent strip on the left
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(1.0))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    add_text_box(slide, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.6),
                 title, size=26, bold=True, color=ACCENT,
                 anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text_box(slide, Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.4),
                     subtitle, size=14, color=GRAY_MID,
                     anchor=MSO_ANCHOR.MIDDLE)
        return Inches(1.2)
    return Inches(0.95)


def add_footer(slide, idx: int, total: int):
    add_text_box(slide, Inches(0.3), Inches(7.1), Inches(8.0), Inches(0.3),
                 "Physical-Mode Activation in Open-Source VLMs",
                 size=9, color=GRAY_LIGHT)
    add_text_box(slide, Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3),
                 f"{idx} / {total}", size=9, color=GRAY_LIGHT,
                 align=PP_ALIGN.RIGHT)


# ---- Slide-builder factories ----

def new_slide(prs):
    blank = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(blank)


# ============================================================================
# Slides
# ============================================================================


def slide_title(prs):
    s = new_slide(prs)
    # Big title block
    add_text_box(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5),
                 "Physical-Mode Activation in Open-Source VLMs",
                 size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.8),
                 "When does an open-source VLM stop seeing a circle as geometry "
                 "and start seeing it as a ball?",
                 size=20, color=GRAY_DARK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Bottom info
    add_text_box(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
                 "Mechanistic interpretability + behavioral psychophysics on Qwen2.5-VL "
                 "/ LLaVA-1.5 / LLaVA-Next / Idefics2 / InternVL3",
                 size=15, color=GRAY_MID, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                 "Review draft  ·  2026-04-26",
                 size=14, color=GRAY_MID, align=PP_ALIGN.CENTER)
    return s


def slide_question(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "The research question",
                      subtitle="From the project plan (references/project.md §1)")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(12.3), Inches(2.0), [
        ("Above what visual-cue threshold does an open-source VLM stop processing "
         "an abstract circle as geometry and start processing it as a physical object?", 18),
        "Two layers of measurement:",
        ("Behavior: PMR (physics-mode reading rate) / GAR (gravity-align rate) / "
         "RC (response consistency).", 14),
        ("Mechanism: vision-encoder probe AUC, LM logit lens, causal patching + "
         "VTI steering + SAE intervention.", 14),
    ])
    # Stim spectrum row at bottom
    figs = ["01_line_blank_none.png", "02_line_ground_none.png",
            "03_shaded_ground_none.png", "04_textured_ground_arrow_shadow.png"]
    labels = ["line/blank/none\n(most abstract)", "line/ground/none\n(+ ground line)",
              "shaded/ground/none\n(3D ball + ground)",
              "textured/ground/\narrow_shadow\n(max cue)"]
    n = len(figs)
    fig_w = Inches(2.7); gap = Inches(0.1)
    total_w = fig_w * n + gap * (n - 1)
    start_x = (SLIDE_W - total_w) / 2
    for i, (f, lbl) in enumerate(zip(figs, labels)):
        x = start_x + (fig_w + gap) * i
        add_figure(s, FIG_DIR / f, x, Inches(4.0), w=fig_w)
        add_text_box(s, x, Inches(6.3), fig_w, Inches(0.6), lbl,
                     size=10, color=GRAY_MID, align=PP_ALIGN.CENTER)
    return s


def slide_tldr(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "TL;DR",
                      subtitle="The three paper-grade headlines")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(12.3), Inches(5.5), [
        ("Open-source VLMs DO read minimal synthetic stimuli as physical, "
         "but the reading is determined at the architecture level "
         "(joint encoder + LM) — not at the encoder representational level.", 17),
        ("Every encoder linearly separates physics-vs-abstract factorial cells "
         "at AUC = 1.0 — encoder discriminability is uniform.", 16),
        ("Behavioral PMR ranges 0.18 → 0.92 across 5 models on identical stim. "
         "The 2-CLIP-point comparison (LLaVA-1.5 0.18 vs LLaVA-Next 0.70) "
         "rules out vision-encoder-family as the sole driver.", 16),
        ("v_L10 — a single LM-layer-10 direction in Qwen2.5-VL — "
         "causally flips behavior under runtime steering AND is "
         "encodable in the image (pixel-space gradient ascent → §4.6).", 16),
        ("Real photographs compress the encoder gap (all 5 models converge "
         "to PMR 0.18-0.67) and halve the label-driven H7 effect — "
         "image-prior dominates label-prior on rich images.", 16),
    ])
    return s


def slide_headline_figure(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "Headline figure: 5-model × 3-stim PMR ladder",
                      subtitle="The 'paper Table 1' visual")
    add_figure(s, FIG_DIR / "session_5model_cross_stim_pmr.png",
               Inches(1.5), y + Inches(0.1), w=Inches(10.3))
    add_caption(s, Inches(1.5), Inches(6.5), Inches(10.3),
                "Mean PMR(_nolabel) ± 95% bootstrap CI per (model × stim source). "
                "Encoder-family split on synthetic stim collapses on real photos.")
    return s


def slide_methods(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "Methodology overview",
                      subtitle="5 sub-tasks, 3 metrics, 3 stim sources, 5 models")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(6.5), Inches(5.5), [
        "Sub-task 1 (PhysCue): behavioral thresholds on a 5-axis factorial.",
        "Sub-task 2 (vision probing): linear probe on encoder activations.",
        "Sub-task 3 (LM logit lens): per-layer probe at visual tokens.",
        "Sub-task 4 (causal): VTI steering (M5a) + SIP/SAE (M5b deferred).",
        "Sub-task 5 (cross-model): 5 VLMs on M8a/M8d/M8c stim.",
    ])
    add_text_box(s, Inches(7.2), y + Inches(0.05), Inches(5.7), Inches(0.4),
                 "Metrics", size=18, bold=True, color=ACCENT)
    add_bullets(s, Inches(7.2), y + Inches(0.5), Inches(5.7), Inches(2.5), [
        ("PMR — fraction of responses that name a physical event "
         "(\"falls\", \"rolls\"). Rule-based scorer, ~5% disagreement vs. "
         "hand-annotation.", 13),
        ("GAR — gravity-align rate; subset of PMR with downward direction.", 13),
        ("RC — response consistency across seeds at T=0.7.", 13),
    ])
    add_text_box(s, Inches(7.2), Inches(4.0), Inches(5.7), Inches(0.4),
                 "Models", size=18, bold=True, color=ACCENT)
    add_bullets(s, Inches(7.2), Inches(4.5), Inches(5.7), Inches(2.5), [
        ("Qwen2.5-VL-7B (SigLIP + Qwen2-7B)", 13),
        ("LLaVA-1.5-7B (CLIP-ViT-L + Vicuna)", 13),
        ("LLaVA-Next-7B (CLIP-ViT-L + Mistral-7B + AnyRes tiling)", 13),
        ("Idefics2-8B (SigLIP-SO400M + Mistral-7B)", 13),
        ("InternVL3-8B (InternViT + InternLM3-8B)", 13),
    ])
    return s


def slide_stim_design(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M2 stimulus design — 5-axis factorial (2880 stim)",
                      subtitle="object_level × bg_level × cue_level × event × seed")
    figs = ["01_line_blank_none.png", "02_line_ground_none.png",
            "03_shaded_ground_none.png", "04_textured_ground_arrow_shadow.png",
            "05_filled_blank_wind.png", "06_textured_blank_none.png"]
    labels = ["line/blank/none", "line/ground/none", "shaded/ground/none",
              "textured/ground/arrow_shadow", "filled/blank/wind",
              "textured/blank/none"]
    fig_w = Inches(2.0); gap = Inches(0.15)
    cols = 3
    for i, (f, lbl) in enumerate(zip(figs, labels)):
        col = i % cols; row = i // cols
        x = Inches(1.4) + (fig_w + gap) * col
        yy = y + Inches(0.1) + (Inches(2.4)) * row
        add_figure(s, FIG_DIR / f, x, yy, w=fig_w)
        add_text_box(s, x, yy + Inches(2.0), fig_w, Inches(0.3), lbl,
                     size=10, color=GRAY_MID, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(8.5), y + Inches(0.5), Inches(4.5), Inches(5.0),
                 "M2 (this slide) is Qwen2.5-VL-only.\n"
                 "Each cell tested with 3 prompts:\n"
                 "• Open: \"What will happen next?\"\n"
                 "• Forced-choice (A/B/C/D)\n"
                 "• Label-free open: \"What do you see?\"\n\n"
                 "PMR is the rule-based scorer over the open response.\n\n"
                 "Cross-model generalization of M2's headlines lives in\n"
                 "M8a / M8d / M8c — see later slides.",
                 size=13, color=GRAY_DARK)
    return s


def slide_m1_pilot(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M1 pilot — first behavioral curve (Qwen2.5-VL, 480 stim)",
                      subtitle="Established H1 partial / H2 strong / H4 / H5 / H6 candidates")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(7.5), Inches(5.5), [
        "Ground effect is the largest single-factor: blank 0.49 → ground 0.85 (+36 pp).",
        "Abstraction endpoints: line 0.58 → textured 0.81 (H1 partial).",
        "Arrow + shadow saturates at PMR = 1.000.",
        "Wind cue invisible to the VLM (PMR ≈ baseline).",
        "Open vs. forced-choice gap = 26 pp → strong language-prior dominance.",
    ])
    add_text_box(s, Inches(8.2), y + Inches(0.1), Inches(4.8), Inches(5.5),
                 "Headline findings:\n\n"
                 "H1: monotone S-curve confirmed (later in M2)\n"
                 "H2: \"ball\" label substantially boosts PMR\n"
                 "H4: open-vs-FC gap = stable signature of\n"
                 "    language-prior ↔ visual-evidence conflict",
                 size=14, color=GRAY_DARK)
    return s


def slide_m2(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M2 MVP-full — 2880-stim factorial (Qwen2.5-VL only)",
                      subtitle="H1 monotone S-curve confirmed; H7 emerged. "
                      "Cross-model generalization of these headlines lives in "
                      "M8a/M8d/M8c (later slides).")
    rows = [
        ["Criterion", "Result"],
        ["Monotone S-curve over object_level (FC)", "0.583 < 0.647 < 0.711 < 0.714"],
        ["Open-vs-FC gap at every level", "22-32 pp"],
        ["cast_shadow alone > none + 20 pp", "+17.5 pp avg"],
        ["RC < 1 cells exist (T>0)", "103/288 (36%)"],
        ["LM activation safetensors captured", "5 layers, bf16"],
    ]
    add_table(s, Inches(0.5), y + Inches(0.1), Inches(8.0), Inches(3.5), rows,
              header_color=ACCENT, font_size=13)
    add_text_box(s, Inches(9.0), y + Inches(0.1), Inches(4.0), Inches(0.4),
                 "New findings", size=18, bold=True, color=ACCENT)
    add_bullets(s, Inches(9.0), y + Inches(0.6), Inches(4.0), Inches(4.0), [
        ("H1 S-curve cleanly confirmed", 13),
        ("H7 emerged: same image with circle/ball/planet → "
         "static / rolls / orbits the Sun", 13),
        ("Open-ended PMR 0.93, abstract-reject 3/1440 — "
         "language-prior dominance reconfirmed", 13),
    ])
    return s


def slide_m3_boomerang(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M3 vision-encoder probing — the 'boomerang'",
                      subtitle="Encoder knows, decoder gates (H-boomerang)")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(6.5), Inches(5.0), [
        "Encoder AUC ≈ 1.0 on every factorial axis from layer 3 onward.",
        ("Behavioral forced-choice PMR ranges 0.28 (no cue) to 0.95 (both) "
         "on the same images.", 14),
        ("Per-object-level encoder AUC ~0.95 constant while behavior "
         "0.58–0.71 → gap is largest at the most abstract end (line, +36 pp).", 14),
        ("→ The information IS there; the LM is gating it. "
         "H-boomerang: 'encoder knows, decoder doesn't'.", 15),
        ("Caveat: programmatic stim makes encoder AUC = 1.0 trivial. "
         "Photo-level extension is M8c.", 13),
    ])
    add_figure(s, FIG_DIR / "encoder_swap_qwen_probe.png",
               Inches(7.5), y + Inches(0.1), w=Inches(5.5))
    return s


def slide_m4_logit_lens(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M4 LM logit lens — per-layer visual-token probe",
                      subtitle="LM AUC plateaus at ~0.95 from L5; switching layer at L20 peak")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(7.5), Inches(5.0), [
        "LM hidden state at visual-token positions linearly separates "
        "physics vs. abstract from L5 (AUC ~0.94, peak L20 = 0.953).",
        ("So the 'forgetting' is NOT during LM forward pass — "
         "the signal is preserved. The bottleneck is at the decoding head.", 14),
        ("M4b (label-free null test, 2026-04-25): paired PMR(ball) − "
         "PMR(_nolabel) = +0.006 ≈ 0; PMR(circle) − PMR(_nolabel) = −0.065. "
         "On Qwen, the language-prior is asymmetric — circle override, "
         "not ball enhancement.", 14),
        ("M4c forced-choice variant: confirms M4b under FC. "
         "LLaVA \"A\" greedy bias surfaces here too.", 14),
    ])
    return s


def slide_m5a_steering(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M5a VTI causal steering — 'physics-mode' direction found",
                      subtitle="L10 α=40 flips 10/10 line/blank/none from "
                      "'D: abstract' to 'B: stays still'")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(7.5), Inches(5.0), [
        ("Recipe: v_L = mean(h_L | PMR=1) − mean(h_L | PMR=0) "
         "from M2-captured activations. v_unit_L = v_L / ||v_L||.", 13),
        ("Test: line/blank/none × 10 seeds × forced-choice with label = circle. "
         "Inject α · v_unit_L at LM layer L over visual-token positions.", 13),
        ("Result: ONLY L10 α=40 flips behavior — 10/10 D→B. "
         "L15/20/25 don't move at the same α.", 14),
        ("Interpretation: v_L10 is a 'physics-mode' linear direction. "
         "Early-mid LM layer is the causal sweet spot — H-locus supported.", 14),
        ("M5a-ext Exp 3 revised this: v_L10 is a regime axis "
         "within physics-mode (+α → kinetic, −α → static), "
         "not a one-way activator.", 14),
    ])
    add_figure(s, FIG_DIR / "01_line_blank_none.png",
               Inches(8.5), y + Inches(0.1), w=Inches(4.5))
    add_caption(s, Inches(8.5), Inches(5.0), Inches(4.5),
                "The stim being steered. v_L10 flips this exact image.")
    return s


def slide_m5a_ext(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M5a-ext — v_L10 is a bidirectional regime axis",
                      subtitle="Exp 1: ceiling. Exp 2: label flip. Exp 3: bidirectional regime.")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(12.3), Inches(5.5), [
        ("Exp 1 (textured/ground/both, near-ceiling baseline): "
         "−α has no effect. Initially framed as 'one-way activator'.", 14),
        ("Exp 2 (line/blank/none × +α=40 with label = ball): "
         "the model switches B (static) → A (rolls / falls). "
         "Label selects regime when the steering direction is active.", 14),
        ("Exp 3 (textured/blank/none, moderate baseline): "
         "−α=40 flips D → B uniformly across (line, textured) × (ball, circle).", 14),
        ("Revised reading: both signs of α activate physics-mode. "
         "Sign selects regime: +α kinetic / −α static. "
         "Baseline D sits BELOW the |α| threshold, not at one end of the axis.", 14),
        ("New hypothesis H-direction-bidirectional supersedes the original "
         "H-regime ('binary object-ness').", 14),
    ])
    return s


def slide_m6_overview(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M6 cross-model — 5 VLMs on the same stim",
                      subtitle="The H-boomerang turns out to be Qwen-scoped; "
                      "the architecture-level reframe takes its place")
    add_figure(s, FIG_DIR / "encoder_chain_5model.png",
               Inches(1.5), y + Inches(0.1), w=Inches(10.3))
    add_caption(s, Inches(1.5), Inches(6.5), Inches(10.3),
                "5-model encoder probe AUC chain (M8a stim). "
                "Non-CLIP cluster: AUC ≥ 0.88. CLIP-ViT-L (LLaVA-1.5 / "
                "LLaVA-Next): AUC 0.73 / 0.77.")
    return s


def slide_m6_r2(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M6 r2 — H-encoder-saturation hypothesis introduced",
                      subtitle="LLaVA-1.5 vision encoder is the bottleneck on its end "
                      "of the chain (AUC ~0.73 vs Qwen ~0.99)")
    fig_w = Inches(6.0)
    add_figure(s, FIG_DIR / "encoder_swap_qwen_probe.png",
               Inches(0.4), y + Inches(0.1), w=fig_w)
    add_figure(s, FIG_DIR / "encoder_swap_llava_probe.png",
               Inches(6.7), y + Inches(0.1), w=fig_w)
    add_caption(s, Inches(0.4), Inches(5.8), fig_w,
                "Qwen2.5-VL SigLIP probe — saturated.")
    add_caption(s, Inches(6.7), Inches(5.8), fig_w,
                "LLaVA-1.5 CLIP-ViT-L probe — unsaturated.")
    return s


def slide_m6_r3_r4(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M6 r3 + r4 — 4-model chain (Idefics2 + InternVL3)",
                      subtitle="3 distinct non-CLIP encoders cluster at AUC ≥ 0.89")
    add_figure(s, FIG_DIR / "encoder_chain_4model.png",
               Inches(1.5), y + Inches(0.1), w=Inches(10.3))
    add_caption(s, Inches(1.5), Inches(6.5), Inches(10.3),
                "4-model AUC chain: Qwen 0.88, LLaVA 0.77, "
                "Idefics2 0.93, InternVL3 0.89 (M8a-stim apples-to-apples). "
                "Only CLIP-ViT-L falls below the saturation cluster.")
    return s


def slide_m6_r6(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M6 r6 — LLaVA-Next 5th model + 2nd CLIP point",
                      subtitle="PMR 0.700 [0.65, 0.74] sits between LLaVA-1.5 "
                      "floor and saturated cluster → encoder family ≠ sole driver")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(7.5), Inches(5.0), [
        ("LLaVA-1.5 (CLIP-ViT-L + Vicuna): PMR(_nolabel) = 0.18.", 14),
        ("LLaVA-Next (CLIP-ViT-L + Mistral + AnyRes tiling): PMR(_nolabel) = 0.70.", 14),
        ("Same encoder family! 0.52-PMR jump rules out vision-encoder family "
         "as sole determinant.", 14),
        ("4-axis confounded (AnyRes / projector / training / LM family) — "
         "can't isolate to LM. Architecture-level reframe.", 14),
        ("H-encoder-saturation now reads: 'joint encoder + LM determines "
         "PMR ceiling', not 'encoder representational capacity alone'.", 14),
    ])
    add_figure(s, FIG_DIR / "encoder_swap_llava_next_probe.png",
               Inches(8.5), y + Inches(0.1), w=Inches(4.5))
    return s


def slide_m8a(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M8a — Non-circle synthetic shapes (5 shapes × 2 models)",
                      subtitle="Strict scoring asymmetry IS the cross-shape "
                      "validation of H-encoder-saturation")
    add_figure(s, FIG_DIR / "m8a_shape_grid.png",
               Inches(0.5), y + Inches(0.1), w=Inches(7.0))
    add_text_box(s, Inches(7.8), y + Inches(0.1), Inches(5.2), Inches(0.4),
                 "Strict pre-registered scoring", size=16, bold=True, color=ACCENT)
    add_bullets(s, Inches(7.8), y + Inches(0.6), Inches(5.2), Inches(5.0), [
        ("Qwen 1/4 PASS, LLaVA 4/4 PASS.", 14),
        ("Saturated encoder → ceiling effect → no behavioral headroom.", 13),
        ("Unsaturated encoder → all four signals (H1, H7, label, gravity) "
         "measurable.", 13),
        ("H1 + H7 are now unsaturated-only.", 13),
        ("Triangle 'wedge' + polygon 'polygon' = label-design weak points.", 13),
    ])
    return s


def slide_m8d(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M8d — Non-ball physical-object categories (car / person / bird)",
                      subtitle="LLaVA 3/3 H7 PASS — strongest cross-category H7 evidence")
    add_figure(s, FIG_DIR / "m8d_full_scene_samples.png",
               Inches(0.5), y + Inches(0.1), w=Inches(7.5))
    add_bullets(s, Inches(8.3), y + Inches(0.1), Inches(4.7), Inches(5.0), [
        ("LLaVA: car +0.525, person +0.138, bird +0.550 on PMR_regime "
         "(physical − abstract).", 13),
        ("Qwen: 0/3 binary H7 (ceiling-flat), but regime distribution "
         "shows figurine 17.5% static / statue 22.5% static.", 13),
        ("New classify_regime keyword classifier — 5.6% hand-annotation error.", 13),
        ("H7 promoted from 'circle-only' to 'cross-category'.", 13),
    ])
    return s


def slide_m8c(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M8c — Real photographs (60 photos × 5 categories)",
                      subtitle="Photos REDUCE Qwen PMR(_nolabel) 18-48 pp; "
                      "encoder gap collapses on photos")
    add_figure(s, FIG_DIR / "m8c_photo_grid.png",
               Inches(0.5), y + Inches(0.1), w=Inches(7.5))
    add_bullets(s, Inches(8.3), y + Inches(0.1), Inches(4.7), Inches(5.0), [
        ("Qwen physical photo PMR drops 18-48 pp vs synthetic — "
         "synthetic-stim minimality is a co-factor of behavioral saturation.", 13),
        ("LLaVA person photo PMR rises +39 pp vs synthetic — "
         "encoder finally recognizes humans.", 13),
        ("All 3 models converge to PMR [0.18, 0.67] on photos.", 13),
        ("LLaVA H7 partially holds on photos (2/4 binary).", 13),
        ("H-encoder-saturation refined: 'encoder representation + "
         "input-context simplicity'.", 13),
    ])
    return s


def slide_sec45(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "§4.5 Cross-encoder swap — Idefics2 SigLIP-SO400M",
                      subtitle="Causal counterfactual at encoder family level")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(7.5), Inches(5.0), [
        ("Idefics2-8B (SigLIP-SO400M + Mistral-7B) on M8a: "
         "PMR(_nolabel) = 0.882. Patterns identically with Qwen on PMR + H7.", 14),
        ("LLaVA (CLIP-ViT-L + Vicuna): 0.175 — outlier.", 14),
        ("Encoder type (SigLIP vs CLIP) drives PMR ceiling regardless of "
         "LM (Qwen2-7B vs Mistral-7B).", 14),
        ("3-point ladder: Qwen 0.838 / LLaVA 0.175 / Idefics2 0.882. "
         "Causally validates H-encoder-saturation at encoder family.", 14),
    ])
    add_figure(s, FIG_DIR / "encoder_swap_heatmap.png",
               Inches(8.3), y + Inches(0.1), w=Inches(4.7))
    return s


def slide_sec46_intro(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "§4.6 — VTI-reverse counterfactual stim (the flagship adversarial)",
                      subtitle="Pixel-space gradient ascent on Qwen2.5-VL pixel_values")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(12.3), Inches(5.0), [
        ("Question: can a small pixel perturbation along v_L10 — "
         "WITHOUT runtime steering — flip Qwen2.5-VL from "
         "\"circle stays static\" to \"circle falls\"?", 16),
        ("Method: gradient ascent on Qwen2.5-VL post-processor "
         "pixel_values (T_patches × 1176, the patch-flattened normalized "
         "representation), maximizing ⟨mean(h_L10[visual]), v_L10⟩. "
         "Adam, lr=1e-2, n_steps=200.", 14),
        ("Configurations on 5 baseline circle stim:", 15),
        ("    bounded ε ∈ {0.05, 0.1, 0.2}, unconstrained, "
         "and 3 random unit directions at ε = 0.1 (matched magnitude).", 13),
    ])
    return s


def slide_sec46_result(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "§4.6 result — 5/5 v_L10 flips at ε = 0.05; 0/15 random",
                      subtitle="Directional specificity falsifies "
                      "'any pixel perturbation flips PMR'")
    rows = [
        ["Config", "n flipped", "Mean final projection"],
        ["bounded ε=0.05 (v_L10)", "5 / 5", "43.7"],
        ["bounded ε=0.10 (v_L10)", "5 / 5", "100.6"],
        ["bounded ε=0.20 (v_L10)", "5 / 5", "125.9"],
        ["unconstrained (v_L10)", "5 / 5", "181.1"],
        ["control random unit dir × 3 @ ε=0.10", "0 / 15", "73-85"],
    ]
    add_table(s, Inches(0.5), y + Inches(0.1), Inches(7.5), Inches(3.5), rows,
              header_color=ACCENT, font_size=14)
    add_bullets(s, Inches(0.5), Inches(5.0), Inches(7.5), Inches(2.0), [
        ("Random projections (~80) ≈ bounded ε=0.1 v_L10 (~101) — "
         "directional specificity, not magnitude, controls the regime flip.", 13),
        ("Pre-registered success criterion was ≥ 3/5 — actual is unambiguous 5/5.", 14),
    ])
    add_figure(s, FIG_DIR / "sec4_6_counterfactual_stim_trajectory.png",
               Inches(8.3), y + Inches(0.1), w=Inches(4.7))
    return s


def slide_sec46_panel(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "§4.6 canonical seed — what the synthesized stim looks like",
                      subtitle="ε = 0.05 introduces visible texture but preserves the "
                      "abstract circle gestalt; introduces NO human-readable physical content")
    add_figure(s, FIG_DIR / "sec4_6_counterfactual_stim_panels.png",
               Inches(0.4), y + Inches(0.1), w=Inches(12.5))
    add_caption(s, Inches(0.4), Inches(6.4), Inches(12.5),
                "baseline → ε=0.05 → ε=0.1 → unconstrained. "
                "The abstract circle gestalt is preserved through all bounded conditions.")
    return s


def slide_m9(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "M9 generalization audit — paper Table 1",
                      subtitle="3 models × 3 stim sources × 5000-iter bootstrap CIs")
    add_figure(s, FIG_DIR / "m9_table1_heatmap.png",
               Inches(0.5), y + Inches(0.1), w=Inches(7.5))
    add_bullets(s, Inches(8.3), y + Inches(0.1), Inches(4.7), Inches(5.0), [
        ("Encoder family causally drives synthetic-stim PMR ceiling: "
         "non-CLIP CIs [0.80, 0.92] vs CLIP [0.14, 0.37] — fully separated.", 13),
        ("Photos compress encoder gap: all 3 → [0.18, 0.67].", 13),
        ("H7 robust only LLaVA-on-synthetic.", 13),
        ("H-LM-modulation flagged but not paper-defensible (Idefics2 M8d "
         "H7 CI just touches 0).", 13),
        ("Replaces PASS/FAIL binarization with bootstrap CI separation.", 13),
    ])
    return s


def slide_sec43(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "§4.3 — Multilingual label switching (Korean / Japanese / Chinese)",
                      subtitle="5 VLMs × 2 non-English languages on M8a circle stim")
    add_figure(s, FIG_DIR / "sec4_3_korean_vs_english_cross_model.png",
               Inches(0.4), y + Inches(0.1), w=Inches(7.5))
    add_bullets(s, Inches(8.0), y + Inches(0.1), Inches(5.0), Inches(5.0), [
        ("Korean (공/원/행성): cross-label ordering preserved 4/5 models.", 14),
        ("LLaVA-1.5 has biggest swing (avg |Δ| = 0.11) — Vicuna's "
         "Korean coverage is weakest.", 13),
        ("Japanese (ボール/円/惑星): tests different mechanism — Qwen "
         "genuinely engages JA (label-echo 85-91%); LLaVA-1.5 translates "
         "kanji to English internally; Idefics2 falls back to Chinese "
         "on 惑星 (24% of responses).", 13),
        ("Scorer extended to KO/JA/CN substring matching (51 → 54 cases).", 13),
    ])
    return s


def slide_sec47(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "§4.7 — Per-axis RC: saturation as decision-stability ceiling",
                      subtitle="cue_level is the dominant decision stabilizer "
                      "for non-CLIP models")
    add_figure(s, FIG_DIR / "sec4_7_rc_per_axis.png",
               Inches(0.4), y + Inches(0.1), w=Inches(7.5))
    add_bullets(s, Inches(8.0), y + Inches(0.1), Inches(5.0), Inches(5.0), [
        ("Saturated models (Qwen 0.84→1.00, Idefics2 0.88→0.99, "
         "InternVL3 0.89→0.98) converge to same PMR call across 5 seeds "
         "when cues fire.", 13),
        ("CLIP-based models (LLaVA-1.5 / LLaVA-Next) retain seed-level "
         "variance even under strong cues.", 13),
        ("Saturation is not just a behavioral PMR ceiling but also a "
         "decision-stability ceiling — separate signature of "
         "H-encoder-saturation.", 13),
    ])
    return s


def slide_sec410(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "§4.10 — Cross-model attention to visual tokens",
                      subtitle="Last-token attention varies architecturally despite "
                      "all 5 LMs receiving 79–98% visual tokens")
    add_figure(s, FIG_DIR / "session_attention_cross_model.png",
               Inches(0.4), y + Inches(0.1), w=Inches(7.5))
    add_bullets(s, Inches(8.0), y + Inches(0.1), Inches(5.0), Inches(5.0), [
        ("Qwen ~17%, LLaVA-1.5 ~7%, Idefics2 ~30%.", 14),
        ("Visual attention peaks at mid-layers in all models.", 14),
        ("Cross-model architecture difference, not encoder difference: "
         "attention allocation is a downstream LM-side property.", 14),
    ])
    return s


def slide_sec411(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "§4.11 — Categorical regime distribution under saturation",
                      subtitle="kinetic / static / abstract / ambiguous fractions "
                      "per (model × category × label_role)")
    add_figure(s, FIG_DIR / "sec4_11_regime_distribution_5model.png",
               Inches(0.4), y + Inches(0.1), w=Inches(8.0))
    add_bullets(s, Inches(8.5), y + Inches(0.1), Inches(4.5), Inches(5.0), [
        ("Granular form of M9 H7 finding.", 13),
        ("InternVL3 person × exotic (statue): PMR drops 0.800 → 0.481, "
         "65% static — strongest single label-driven static commit.", 13),
        ("LLaVA-1.5 most regime-discriminative.", 13),
        ("Categorical view reveals the KIND of commitment, not just whether "
         "the model commits.", 13),
    ])
    return s


def slide_hypotheses(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "Hypothesis scorecard — final state",
                      subtitle="Original H1-H7 + named H- hypotheses derived during the work")
    rows = [
        ["ID", "Status", "Key evidence"],
        ["H1 (S-curve)", "supported, unsaturated-only", "M2 (Qwen, compressed) + M6 r1 (LLaVA clean) + M8a (LLaVA 4/5)"],
        ["H2 (label prior)", "validated, encoder-anchored", "M4b (Qwen) + M6 r1 (LLaVA) + M6 r2a (InternVL3)"],
        ["H4 (open vs FC)", "supported (Qwen-only at M2)", "M2: 22-32 pp gap; cross-model untested"],
        ["H5 (ground)", "mixed (Qwen-only)", "M2: bg +21 vs obj +9; cross-model untested"],
        ["H6 (cast shadow)", "supported, revised (Qwen-only)", "M2: arrow alone also saturates"],
        ["H7 (label selects regime)", "validated, cross-category", "M2 (Qwen GAR) + M5a-ext + M6 r1+r2a (circle) + M8d LLaVA 3/3"],
        ["H-boomerang", "Qwen-scoped", "M3 (Qwen): AUC 1.0; refuted on LLaVA-1.5"],
        ["H-encoder-saturation", "architecture-level confirmed", "M9 5-model bootstrap (3 stim)"],
        ["H-LM-modulation", "suggested only", "M9 Idefics2 M8d H7 CI touches 0"],
        ["H-locus (mid-LM)", "supported (Qwen-only)", "M5a (Qwen): L10 α=40 only"],
        ["H-direction-bidirectional", "supported (Qwen-only)", "M5a-ext Exp 3 (Qwen)"],
        ["H-direction-specificity", "supported (Qwen-only, §4.6)", "5/5 v_L10 vs 0/15 random (Qwen)"],
        ["H-shortcut", "supported (Qwen-only, §4.6)", "Pixel-encodable on Qwen"],
    ]
    add_table(s, Inches(0.4), y + Inches(0.05), Inches(12.5), Inches(5.8), rows,
              header_color=ACCENT, font_size=11)
    return s


def slide_limitations(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "Limitations carried forward",
                      subtitle="What's still open / honest caveats")
    add_bullets(s, Inches(0.5), y + Inches(0.1), Inches(12.3), Inches(5.5), [
        "Single-task: 'next-state-prediction prompt' is the only behavioral readout.",
        ("Programmatic stim makes encoder AUC = 1.0 trivial. M8c photos "
         "partially address this; richer photo distributions still open.", 14),
        ("v_L10 is a 1-d axis from PCA over a labeled stim distribution. "
         "Multi-direction decomposition (§4.6 limitation) is open.", 14),
        ("§4.6 cross-model: each model needs its own M5a-style v_L10 "
         "computation (not yet done for LLaVA / Idefics2 / InternVL3).", 14),
        ("M5b (SIP / activation patching / SAE) is the major mechanism gap.", 14),
        "M7 Prolific human baseline not yet collected.",
        ("H-LM-modulation suggested only; clean LM-only counterfactual "
         "would need a same-encoder LM-swap.", 14),
    ])
    return s


def slide_open_backlog(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "Open backlog & future work",
                      subtitle="Roadmap §3 status as of 2026-04-26")
    rows = [
        ["Item", "Type", "Status"],
        ["M5b — SIP / activation patching / SAE", "Major mechanism", "Optional / next paper-section gap"],
        ["§4.4 — Michotte 2-frame causality", "New investigation", "Needs 2-image prompt infra"],
        ["§4.6 cross-model", "Extension", "Needs per-model v_L10 (M5a-style on each)"],
        ["§4.8 — PMR scaling (Qwen 32B / 72B)", "Scale check", "Needs new large-model loads"],
        ["M7 — Prolific human baseline", "External", "20 raters × 50 stim, ~$200"],
        ["M7 — paper draft (EMNLP / NeurIPS)", "Writing", "Outline ready (this PPT)"],
    ]
    add_table(s, Inches(0.4), y + Inches(0.1), Inches(12.5), Inches(5.0), rows,
              header_color=ACCENT, font_size=12)
    return s


def slide_summary(prs):
    s = new_slide(prs)
    y = add_title_bar(s, "Summary",
                      subtitle="Three paper-grade contributions")
    add_bullets(s, Inches(0.5), y + Inches(0.2), Inches(12.3), Inches(5.5), [
        ("Architecture-level reframe of physics-mode reading: "
         "5-model × 3-stim bootstrap-CI evidence that encoder family alone "
         "doesn't determine PMR ceiling. The 2-CLIP-point insight (LLaVA-1.5 "
         "0.18 vs LLaVA-Next 0.70) is the cleanest disconfirmer.", 16),
        ("Causal localization of physics-mode reading: M5a flips behavior at "
         "exactly LM L10 with α=40. M5a-ext shows v_L10 is a "
         "regime axis (+α kinetic, −α static), not a binary toggle.", 16),
        ("v_L10 is encodable in the image: §4.6 pixel-space gradient ascent "
         "produces 5/5 PMR flips at ε=0.05 with 0/15 random-direction controls. "
         "The shortcut path can be 'spelled' in pixels.", 16),
        ("Plus: cross-shape (M8a) + cross-category (M8d) + photo-collapse (M8c) "
         "+ multilingual (§4.3) + decision-stability ceiling (§4.7) + "
         "regime distribution (§4.11) — all consistent with the architecture "
         "reframe.", 14),
    ])
    return s


def slide_thanks(prs):
    s = new_slide(prs)
    add_text_box(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2),
                 "Thank you", size=60, bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.8),
                 "Questions / discussion / next-step prioritization",
                 size=22, color=GRAY_DARK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
                 "Repo: github.com/namam3gy/physical-mode-activation",
                 size=14, color=GRAY_MID, align=PP_ALIGN.CENTER)
    return s


# ---- Table helper ----

def add_table(slide, x, y, w, h, rows: list[list[str]], *,
              header_color=ACCENT, font_size=12):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = table_shape.table
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    if i == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    RGBColor(0xF6, 0xF8, 0xFC) if i % 2 == 0
                    else RGBColor(0xFF, 0xFF, 0xFF))
    return table_shape


# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        ("Title", slide_title),
        ("Question", slide_question),
        ("TL;DR", slide_tldr),
        ("Headline figure", slide_headline_figure),
        ("Methodology", slide_methods),
        ("Stim design", slide_stim_design),
        ("M1 pilot", slide_m1_pilot),
        ("M2", slide_m2),
        ("M3 boomerang", slide_m3_boomerang),
        ("M4 logit lens", slide_m4_logit_lens),
        ("M5a steering", slide_m5a_steering),
        ("M5a-ext", slide_m5a_ext),
        ("M6 overview", slide_m6_overview),
        ("M6 r2", slide_m6_r2),
        ("M6 r3 + r4", slide_m6_r3_r4),
        ("M6 r6", slide_m6_r6),
        ("M8a", slide_m8a),
        ("M8d", slide_m8d),
        ("M8c", slide_m8c),
        ("§4.5", slide_sec45),
        ("§4.6 intro", slide_sec46_intro),
        ("§4.6 result", slide_sec46_result),
        ("§4.6 panels", slide_sec46_panel),
        ("M9", slide_m9),
        ("§4.3", slide_sec43),
        ("§4.7", slide_sec47),
        ("§4.10", slide_sec410),
        ("§4.11", slide_sec411),
        ("Hypotheses", slide_hypotheses),
        ("Limitations", slide_limitations),
        ("Open backlog", slide_open_backlog),
        ("Summary", slide_summary),
        ("Thanks", slide_thanks),
    ]

    total = len(builders)
    for i, (name, fn) in enumerate(builders, 1):
        s = fn(prs)
        if i > 1 and i < total:  # no footer on title or thanks
            add_footer(s, i, total)
        print(f"  [{i:2d}/{total}] {name}")

    out_dir = PROJECT_ROOT / "docs" / "review_ppt"
    out_dir.mkdir(parents=True, exist_ok=True)
    ts = time.strftime("%Y%m%d-%H%M%S")
    out_path = out_dir / f"physical_mode_review_{ts}.pptx"
    prs.save(str(out_path))
    print(f"\nSaved: {out_path}")
    # Also save a stable copy at a fixed name for easy linking.
    stable = out_dir / "physical_mode_review_latest.pptx"
    prs.save(str(stable))
    print(f"Stable: {stable}")


if __name__ == "__main__":
    main()
